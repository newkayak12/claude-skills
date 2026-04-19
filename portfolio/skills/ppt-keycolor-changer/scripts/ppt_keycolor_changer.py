"""
PPT Key Color Changer
Replace specific colors across all slides in a PPTX file.

Usage:
    python ppt_keycolor_changer.py \
        --input file.pptx \
        --mapping '{"003087": "0064FF"}' \
        --tolerance 5 \
        --output file_recolored.pptx
"""

import argparse
import json
import sys
from collections import defaultdict
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.util import Pt
except ImportError:
    print("ERROR: python-pptx is required. Run: pip install python-pptx")
    sys.exit(1)


def hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    h = hex_str.lstrip("#").upper()
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def color_distance(c1: tuple, c2: tuple) -> float:
    return sum((a - b) ** 2 for a, b in zip(c1, c2)) ** 0.5


def find_target(rgb: tuple, mapping_rgb: dict, tolerance: int) -> str | None:
    """Return new hex if rgb matches any source color within tolerance."""
    for src_rgb, dst_hex in mapping_rgb.items():
        if color_distance(rgb, src_rgb) <= tolerance:
            return dst_hex
    return None


def try_replace_color(color_obj, mapping_rgb: dict, tolerance: int) -> str | None:
    """Try to read and replace a color object. Returns dst_hex if replaced, else None."""
    try:
        if color_obj.type is None:
            return None
        rgb = color_obj.rgb
        current = (rgb.r, rgb.g, rgb.b)
        dst_hex = find_target(current, mapping_rgb, tolerance)
        if dst_hex:
            r, g, b = hex_to_rgb(dst_hex)
            color_obj.rgb = RGBColor(r, g, b)
            return dst_hex
    except Exception:
        pass
    return None


def process_fill(fill, mapping_rgb, tolerance):
    """Returns dst_hex if replaced."""
    try:
        return try_replace_color(fill.fore_color, mapping_rgb, tolerance)
    except Exception:
        return None


def process_shape(shape, mapping_rgb, tolerance, counts):
    # Fill
    try:
        result = process_fill(shape.fill, mapping_rgb, tolerance)
        if result:
            counts[("fill", result)] += 1
    except Exception:
        pass

    # Line color
    try:
        result = try_replace_color(shape.line.color, mapping_rgb, tolerance)
        if result:
            counts[("line", result)] += 1
    except Exception:
        pass

    # Text
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    result = try_replace_color(run.font.color, mapping_rgb, tolerance)
                    if result:
                        counts[("text", result)] += 1
    except Exception:
        pass

    # Table cells
    try:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    result = process_fill(cell.fill, mapping_rgb, tolerance)
                    if result:
                        counts[("table_fill", result)] += 1
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            result = try_replace_color(run.font.color, mapping_rgb, tolerance)
                            if result:
                                counts[("table_text", result)] += 1
    except Exception:
        pass

    # Chart series
    try:
        if shape.has_chart:
            chart = shape.chart
            for series in chart.series:
                try:
                    result = process_fill(series.format.fill, mapping_rgb, tolerance)
                    if result:
                        counts[("chart_series", result)] += 1
                except Exception:
                    pass
    except Exception:
        pass

    # Group shapes (recursive)
    try:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for child in shape.shapes:
                process_shape(child, mapping_rgb, tolerance, counts)
    except Exception:
        pass


def process_slide_background(slide, mapping_rgb, tolerance, counts):
    try:
        bg = slide.background
        result = process_fill(bg.fill, mapping_rgb, tolerance)
        if result:
            counts[("background", result)] += 1
    except Exception:
        pass


def run(input_path: str, mapping: dict, tolerance: int, output_path: str | None):
    prs = Presentation(input_path)

    # Normalize mapping to uppercase hex without #
    mapping_clean = {k.lstrip("#").upper(): v.lstrip("#").upper() for k, v in mapping.items()}
    mapping_rgb = {hex_to_rgb(src): dst for src, dst in mapping_clean.items()}

    total_counts: dict[tuple, int] = defaultdict(int)

    for slide_idx, slide in enumerate(prs.slides, 1):
        process_slide_background(slide, mapping_rgb, tolerance, total_counts)
        for shape in slide.shapes:
            process_shape(shape, mapping_rgb, tolerance, total_counts)

    if output_path is None:
        p = Path(input_path)
        output_path = str(p.parent / f"{p.stem}_recolored{p.suffix}")

    prs.save(output_path)

    # Print summary
    print(f"\nSaved: {output_path}\n")
    if not total_counts:
        print("WARNING: No colors were replaced.")
        print("  → The specified colors may not exist in the file, or may be theme colors.")
        print("  → Try increasing --tolerance (e.g. --tolerance 10).")
        return

    # Group by dst color
    by_dst: dict[str, dict[str, int]] = defaultdict(lambda: defaultdict(int))
    for (elem_type, dst_hex), count in total_counts.items():
        by_dst[dst_hex][elem_type] += count

    for dst_hex, elem_counts in by_dst.items():
        total = sum(elem_counts.values())
        detail = ", ".join(f"{k}: {v}" for k, v in sorted(elem_counts.items()))
        # Find which src mapped to this dst
        src_hex = next((s for s, d in mapping_clean.items() if d == dst_hex), "?")
        print(f"  #{src_hex} → #{dst_hex}: {total} replacements ({detail})")

    print()


def main():
    parser = argparse.ArgumentParser(description="Replace colors in a PPTX file.")
    parser.add_argument("--input", required=True, help="Path to source .pptx file")
    parser.add_argument("--mapping", required=True, help='JSON: {"OLD_HEX": "NEW_HEX", ...}')
    parser.add_argument("--tolerance", type=int, default=0, help="Color match tolerance 0–30 (default: 0)")
    parser.add_argument("--output", default=None, help="Output path (default: <name>_recolored.pptx)")
    args = parser.parse_args()

    try:
        mapping = json.loads(args.mapping)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid --mapping JSON: {e}")
        sys.exit(1)

    run(args.input, mapping, args.tolerance, args.output)


if __name__ == "__main__":
    main()
